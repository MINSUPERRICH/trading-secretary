import streamlit as st
import pandas as pd
import plotly.express as px
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

# --- 1. INITIALIZATION ---
if "user_db" not in st.session_state:
    st.session_state["user_db"] = {"rich": {"password": "777", "role": "admin"}, "staff": {"password": "123", "role": "team"}}
if "password_correct" not in st.session_state:
    st.session_state["password_correct"] = False
if "master_df" not in st.session_state:
    st.session_state["master_df"] = None

st.set_page_config(page_title="Executive BI Hub", layout="wide")

# --- 2. PPT ENGINE ---
def generate_advanced_ppt(fig, title, insight):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Performance Visualization"
    img_bytes = fig.to_image(format="png")
    slide2.shapes.add_picture(BytesIO(img_bytes), Inches(1), Inches(1.5), width=Inches(8))
    
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Business Insights"
    slide3.placeholders[1].text = insight
    
    ppt_out = BytesIO()
    prs.save(ppt_out)
    return ppt_out.getvalue()

# --- 3. LOGIN SYSTEM ---
def check_password():
    if not st.session_state["password_correct"]:
        st.title("🔒 Business Intelligence Login")
        u = st.text_input("Username")
        p = st.text_input("Password", type="password")
        if st.button("Log In"):
            if u in st.session_state["user_db"] and st.session_state["user_db"][u]["password"] == p:
                st.session_state.update({"password_correct": True, "role": st.session_state["user_db"][u]["role"], "user": u})
                st.rerun()
            else: st.error("Invalid credentials")
        return False
    return True

# --- 4. MAIN APP ---
if check_password():
    st.sidebar.title(f"Welcome, {st.session_state['user'].title()}")
    
    tab_dash, tab_pivot, tab_merge, tab_ppt, tab_settings = st.tabs([
        "📊 Dashboard", "🧮 Pivot Table", "🔗 Link & Match", "🎬 PPT Designer", "⚙️ Settings"
    ])

    uploaded_files = st.sidebar.file_uploader("Upload Data Files", type=['csv', 'xlsx'], accept_multiple_files=True)
    
    fig = None 

    if uploaded_files:
        data_dict = {}
        for f in uploaded_files:
            data_dict[f.name] = pd.read_excel(f) if f.name.endswith('xlsx') else pd.read_csv(f)
        
        if st.session_state["master_df"] is None or len(uploaded_files) == 1:
            st.session_state["master_df"] = list(data_dict.values())[0]
            
        working_df = st.session_state["master_df"]

        # --- AUTO-RENAME MACD & DAILY COLUMNS ---
        rename_dict = {
            "4H_Signal_Now": "4H_MACD_Signal_Now",
            "4H_Signal_Prev": "4H_MACD_Signal_Prev",
            "4H Signal": "4H_MACD_Trend",
            "1D_Signal_Now": "D_MACD_Signal_Now",
            "1D_Signal_Prev": "D_MACD_Signal_Prev",
            "1D Signal": "D_MACD_Trend",
            "Signal_Now": "D_MACD_Signal_Now",
            "Signal_Prev": "D_MACD_Signal_Prev",
            "MACD_Now": "D_MACD_Now",
            "MACD_Prev": "D_MACD_Prev"
        }
        working_df = working_df.rename(columns=rename_dict)
        working_df.columns = working_df.columns.str.replace("Daily", "D", regex=False)
        st.session_state["master_df"] = working_df
        # ----------------------------------------

        # --- DASHBOARD TAB ---
        with tab_dash:
            st.header("Visual Performance")
            num_cols = working_df.select_dtypes('number').columns.tolist()
            if num_cols:
                col = st.sidebar.selectbox("Select KPI", num_cols)
                grp = st.sidebar.selectbox("Group By", working_df.columns)
                
                fig = px.bar(working_df.groupby(grp)[col].sum().reset_index(), x=grp, y=col, color_discrete_sequence=['#4eb8b8'])
                fig.update_layout(template="plotly_white")
                st.plotly_chart(fig, use_container_width=True)

        # --- PIVOT TABLE TAB ---
        with tab_pivot:
            st.header("Data Pivot Summary")
            if num_cols:
                st.dataframe(working_df.groupby(grp)[col].sum(), use_container_width=True)

        # --- LINK & MATCH (VLOOKUP) TAB ---
        with tab_merge:
            st.header("🔗 Link & Match (VLOOKUP)")
            
            if len(data_dict) < 2:
                st.info("Upload at least TWO files in the sidebar to merge them.")
            else:
                col1, col2 = st.columns(2)
                with col1:
                    base_file = st.selectbox("1. Main File", list(data_dict.keys()))
                with col2:
                    merge_file = st.selectbox("2. File to Link", [f for f in data_dict.keys() if f != base_file])
                
                m_col1, m_col2 = st.columns(2)
                with m_col1:
                    left_on = st.selectbox(f"Key in {base_file}", data_dict[base_file].columns)
                with m_col2:
                    right_on = st.selectbox(f"Key in {merge_file}", data_dict[merge_file].columns)
                
                if st.button("Run VLOOKUP Match"):
                    new_merged_df = pd.merge(data_dict[base_file], data_dict[merge_file], left_on=left_on, right_on=right_on, how='left')
                    st.session_state["master_df"] = new_merged_df
                    st.success("Files Linked Successfully!")

            if st.session_state["master_df"] is not None:
                st.write("**👁️ Preview of Linked Data (First 50 Rows):**")
                st.dataframe(st.session_state["master_df"].head(50), use_container_width=True)

            st.divider()
            st.subheader("🔍 Deep Dive & Filter")
            
            use_date_filter = st.toggle("📅 Enable Date Filter")
            selected_dates = None
            date_col = None
            
            if use_date_filter:
                date_col = st.selectbox("Select Date Column", working_df.columns)
                temp_dates = pd.to_datetime(working_df[date_col], errors='coerce')
                
                if temp_dates.notna().any():
                    min_date = temp_dates.min().date()
                    max_date = temp_dates.max().date()
                    selected_dates = st.date_input("Select Date Range", value=(min_date, max_date), min_value=min_date, max_value=max_date)
                else:
                    st.warning("⚠️ The selected column doesn't seem to contain valid dates.")

            with st.expander("💡 Instruction Cheat Sheet (Click to open)"):
                st.markdown("""
                **How to write filter instructions:**
                * **Numbers:** Use `>`, `<`, `==` (equals). Example: `close > 150`
                * **Text/Words:** Use `==` and put the word in double quotes. Example: `ticker == "AAPL"`
                * **MACD Cross Rule:** `4H_MACD_Now > 4H_MACD_Signal_Now`
                """)
            
            filter_query = st.text_input("Filter Instruction (Optional if using Dates):")
            
            if st.button("Apply Filters"):
                try:
                    filtered_df = working_df.copy()
                    
                    if use_date_filter and selected_dates and len(selected_dates) == 2:
                        start_date, end_date = selected_dates
                        filtered_df[date_col] = pd.to_datetime(filtered_df[date_col], errors='coerce')
                        filtered_df = filtered_df[(filtered_df[date_col].dt.date >= start_date) & (filtered_df[date_col].dt.date <= end_date)]
                    
                    if filter_query:
                        filtered_df = filtered_df.query(filter_query)
                        
                    st.success(f"Found {len(filtered_df)} matching rows!")
                    
                    filter_num_cols = filtered_df.select_dtypes('number').columns.tolist()
                    if filter_num_cols and not filtered_df.empty:
                        stat_kpi = col if 'col' in locals() and col in filter_num_cols else filter_num_cols[0]
                        st.write(f"**Quick Stats for: {stat_kpi}**")
                        s_col1, s_col2, s_col3 = st.columns(3)
                        s_col1.metric("Total (Sum)", f"{filtered_df[stat_kpi].sum():,.2f}")
                        s_col2.metric("Average", f"{filtered_df[stat_kpi].mean():,.2f}")
                        s_col3.metric("Max Value", f"{filtered_df[stat_kpi].max():,.2f}")
                    
                    st.dataframe(filtered_df, use_container_width=True)
                    
                    st.write("**Export Filtered Data:**")
                    dl_col1, dl_col2 = st.columns(2)
                    csv_data = filtered_df.to_csv(index=False).encode('utf-8')
                    
                    excel_buffer = BytesIO()
                    with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
                        df_to_export = filtered_df.copy()
                        for c in df_to_export.select_dtypes(include=['datetimetz']).columns:
                            df_to_export[c] = df_to_export[c].dt.tz_localize(None)
                        df_to_export.to_excel(writer, index=False, sheet_name='Filtered Data')
                    excel_data = excel_buffer.getvalue()
                    
                    with dl_col1:
                        st.download_button("📥 Download CSV", data=csv_data, file_name="Filtered_Data.csv", mime="text/csv")
                    with dl_col2:
                        st.download_button("📊 Download Excel", data=excel_data, file_name="Filtered_Data.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                    
                except Exception as e:
                    st.error("Could not apply filter. Please check your spelling or formatting against the Cheat Sheet.")
            
            if st.session_state["master_df"] is not None and not filter_query and not use_date_filter:
                st.write("**Export Full Merged Data:**")
                f_csv = st.session_state["master_df"].to_csv(index=False).encode('utf-8')
                
                f_excel_buffer = BytesIO()
                with pd.ExcelWriter(f_excel_buffer, engine='xlsxwriter') as writer:
                    df_to_export_full = st.session_state["master_df"].copy()
                    for c in df_to_export_full.select_dtypes(include=['datetimetz']).columns:
                        df_to_export_full[c] = df_to_export_full[c].dt.tz_localize(None)
                    df_to_export_full.to_excel(writer, index=False, sheet_name='Merged Data')
                f_excel = f_excel_buffer.getvalue()
                
                f_col1, f_col2 = st.columns(2)
                with f_col1:
                    st.download_button("📥 Download Full CSV", f_csv, "Full_Merged_Dataset.csv")
                with f_col2:
                    st.download_button("📊 Download Full Excel", f_excel, "Full_Merged_Dataset.xlsx")

        # --- PPT DESIGNER TAB ---
        with tab_ppt:
            st.header("PPT Presentation Designer")
            if fig is not None:
                p_title = st.text_input("Slide Title", "Business Review")
                p_insight = st.text_area("Observations", "Add your meeting notes here.")
                if st.button("Generate Presentation"):
                    ppt_data = generate_advanced_ppt(fig, p_title, p_insight)
                    st.download_button("📥 Download PowerPoint", ppt_data, "Meeting_Deck.pptx")
            else:
                st.warning("Please configure your chart in the 'Dashboard' tab first.")
    else:
        st.info("Please upload your data in the sidebar to begin.")

    if st.sidebar.button("Logout"):
        st.session_state.clear()
        st.rerun()
